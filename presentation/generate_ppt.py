from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUTPUT = Path(
    "/Users/wzm/学习/工作/喷漆/code/repo_zip_download_20260326_0153/extracted/--main/presentation/Smart_Parking_MQTT_Presentation_EN.pptx"
)

BG = RGBColor(247, 249, 252)
NAVY = RGBColor(25, 44, 76)
BLUE = RGBColor(35, 91, 166)
TEAL = RGBColor(42, 143, 136)
ORANGE = RGBColor(222, 133, 57)
LIGHT = RGBColor(228, 236, 247)
TEXT = RGBColor(46, 56, 70)
WHITE = RGBColor(255, 255, 255)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_title(slide, title, subtitle=""):
    add_top_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(12.0), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Aptos"
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(13)
        run2.font.color.rgb = TEXT
        run2.font.name = "Aptos"


def add_textbox(slide, left, top, width, height, paragraphs, fill=None, color=TEXT):
    if fill is not None:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        box = slide.shapes.add_textbox(
            Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.36), Inches(height - 0.24)
        )
    else:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for para in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if isinstance(para, tuple):
            text, size, bold = para
        else:
            text, size, bold = para, 17, False
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color if fill is None else (WHITE if fill != LIGHT else NAVY)
        p.font.name = "Aptos"
        p.space_after = Pt(6)
    return box


def add_bullet_box(slide, left, top, width, height, title, items, fill=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LIGHT

    box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.36), Inches(height - 0.24))
    tf = box.text_frame
    tf.word_wrap = True
    head = tf.paragraphs[0]
    head.text = title
    head.font.size = Pt(19)
    head.font.bold = True
    head.font.color.rgb = NAVY
    head.font.name = "Aptos"
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT
        p.font.name = "Aptos"
        p.space_after = Pt(4)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_top_bar(slide)
hero = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.1), Inches(11.7), Inches(5.2)
)
hero.fill.solid()
hero.fill.fore_color.rgb = WHITE
hero.line.color.rgb = LIGHT
title_box = slide.shapes.add_textbox(Inches(1.15), Inches(1.65), Inches(10.5), Inches(2.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Smart Parking Prototype Based on MQTT"
r.font.size = Pt(30)
r.font.bold = True
r.font.color.rgb = NAVY
r.font.name = "Aptos"
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "Week 4-5 Coursework Presentation"
r2.font.size = Pt(18)
r2.font.color.rgb = BLUE
r2.font.name = "Aptos"
p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = (
    "This project presents a simple Internet of Things prototype for a smart parking lot. "
    "The system uses MQTT to connect parking sensors, a display service, a statistics service, "
    "and an entry gate controller in one shared message-driven workflow."
)
r3.font.size = Pt(19)
r3.font.color.rgb = TEXT
r3.font.name = "Aptos"
chip = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.15), Inches(4.95), Inches(2.55), Inches(0.6)
)
chip.fill.solid()
chip.fill.fore_color.rgb = TEAL
chip.line.fill.background()
chip.text_frame.text = "Python + MQTT"
chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
chip.text_frame.paragraphs[0].font.size = Pt(16)
chip.text_frame.paragraphs[0].font.bold = True
chip.text_frame.paragraphs[0].font.color.rgb = WHITE
chip.text_frame.paragraphs[0].font.name = "Aptos"
footer = slide.shapes.add_textbox(Inches(1.15), Inches(5.75), Inches(10.2), Inches(0.5))
footer.text_frame.text = (
    "The aim of this prototype is not to build a full commercial platform, but to demonstrate a clear and working MQTT communication architecture."
)
footer.text_frame.paragraphs[0].font.size = Pt(15)
footer.text_frame.paragraphs[0].font.color.rgb = TEXT
footer.text_frame.paragraphs[0].font.name = "Aptos"


# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "1. Background and Problem Statement",
    "A parking lot needs fast status updates and better coordination between physical devices.",
)
add_bullet_box(
    slide,
    0.85,
    1.75,
    5.7,
    4.8,
    "Why this topic matters",
    [
        "Drivers often waste time searching for available spaces, especially when the parking lot is busy.",
        "Traditional parking management is slow because information is updated manually or displayed too late.",
        "A message-based system can provide real-time communication between devices and services.",
        "MQTT is suitable for this scenario because it is lightweight, simple, and widely used in IoT projects.",
    ],
)
add_bullet_box(
    slide,
    6.8,
    1.75,
    5.7,
    4.8,
    "What our prototype tries to solve",
    [
        "Detect whether parking spaces are free or occupied through simulated sensors.",
        "Publish those updates to a broker so that multiple services can receive the same data immediately.",
        "Display the latest parking information in a simple monitoring dashboard.",
        "Use the number of free spaces to decide whether the entry gate should open for a new car.",
    ],
)


# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "2. System Architecture",
    "The design follows a publish / subscribe model with one public MQTT broker in the center.",
)
for label, x, y, color in [
    ("Parking Sensors\npublisher_sensor.py", 0.9, 2.2, TEAL),
    ("Vehicle Request Generator\npublisher_vehicle_request.py", 0.9, 4.45, ORANGE),
    ("MQTT Broker\nbroker.hivemq.com", 5.4, 3.2, NAVY),
    ("Display Service\nsubscriber_display.py", 10.05, 1.55, TEAL),
    ("Statistics Service\nsubscriber_stats.py", 10.05, 3.2, ORANGE),
    ("Gate Controller\nsubscriber_gateway.py", 10.05, 4.85, TEAL),
]:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.35), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = label
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for p in tf.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Aptos"
for arrow, x, y in [
    ("→", 3.7, 2.45),
    ("→", 3.7, 4.7),
    ("→", 8.75, 1.9),
    ("→", 8.75, 3.55),
    ("→", 8.75, 5.2),
]:
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.4), Inches(0.4))
    p = t.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = arrow
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = NAVY
arch_text = (
    "The architecture separates data producers from data consumers. Parking sensors publish slot status messages. "
    "A second publisher simulates vehicle arrival requests at the entrance. The display service subscribes to all parking updates and shows the latest state. "
    "The statistics service calculates how many spaces are free and publishes a summary. The gate controller listens to both the vehicle request topic and the statistics topic, "
    "then decides whether to allow a car to enter. This structure keeps the system modular and easy to extend."
)
add_textbox(slide, 0.95, 6.0, 11.5, 0.8, [(arch_text, 14, False)], fill=LIGHT, color=NAVY)


# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "3. MQTT Topics and Data Flow",
    "Each topic has a clear responsibility so that the message flow remains easy to understand.",
)
table = slide.shapes.add_table(6, 3, Inches(0.9), Inches(1.8), Inches(11.5), Inches(3.6)).table
headers = ["Topic", "Publisher", "Purpose"]
for c, text in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Aptos"
rows = [
    ("parking/lot/A01/status", "Parking sensors", "Publishes the current state of one specific parking slot."),
    ("parking/lot/+/status", "All parking sensors", "Allows subscribers to receive updates from all parking slots."),
    ("parking/stats/summary", "Statistics service", "Sends the total number of free and occupied slots."),
    ("parking/gate/request", "Vehicle request generator", "Represents a car asking to enter the parking area."),
    ("parking/gate/entry", "Gate controller", "Returns the final decision such as open or wait."),
]
for r_i, row in enumerate(rows, start=1):
    for c_i, value in enumerate(row):
        cell = table.cell(r_i, c_i)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT
        p.font.name = "Aptos"
flow_text = (
    "The overall data flow is straightforward. Sensor messages are published first, then the statistics service aggregates them and produces a summary. "
    "When a vehicle request arrives, the gate controller checks the latest summary and sends an entry decision. "
    "Because all modules communicate through the broker, the system remains loosely coupled and more flexible than a direct device-to-device design."
)
add_textbox(slide, 0.95, 5.7, 11.45, 1.0, [(flow_text, 14, False)], fill=LIGHT, color=NAVY)


# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "4. Implementation Details",
    "The codebase is small, but each Python module has a distinct role in the prototype.",
)
add_bullet_box(
    slide,
    0.85,
    1.8,
    4.0,
    4.9,
    "Publishers",
    [
        "publisher_sensor.py simulates five parking slots and publishes free or occupied states at random time intervals.",
        "publisher_vehicle_request.py simulates a car arriving at the entrance and publishes a structured request in JSON format.",
        "Both publishers use the same public broker and send messages continuously to represent a live IoT environment.",
    ],
)
add_bullet_box(
    slide,
    4.95,
    1.8,
    4.0,
    4.9,
    "Subscribers",
    [
        "subscriber_display.py receives updates from all parking slots and prints a real-time text dashboard.",
        "subscriber_stats.py stores the latest state of each slot and publishes a summary containing free, occupied, and unknown counts.",
        "subscriber_gateway.py listens to the summary and the vehicle request topic, then decides whether to open the gate or keep it closed.",
    ],
)
add_bullet_box(
    slide,
    9.05,
    1.8,
    3.45,
    4.9,
    "Small improvements added",
    [
        "A new vehicle request module was added so that the gate logic depends on actual MQTT events.",
        "The gate controller was upgraded from a timer-based demo to a decision-based controller.",
        "The README and requirements file were rewritten to make the project easier to run and present.",
    ],
)


# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "5. Testing and Demonstration Results",
    "The prototype was tested by running the publishers and subscribers at the same time.",
)
left_text = (
    "During testing, the statistics service successfully received parking slot updates and recalculated the number of free spaces. "
    "The gate controller also received those summaries and reacted to vehicle requests in real time. "
    "When at least one space was available, the controller published an open decision. "
    "If no free space was available at that moment, the controller returned a wait decision instead."
)
add_textbox(slide, 0.9, 1.8, 5.5, 2.1, [("Observed behavior", 19, True), (left_text, 15, False)], fill=WHITE)
log_box = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.8), Inches(5.7), Inches(3.25)
)
log_box.fill.solid()
log_box.fill.fore_color.rgb = NAVY
log_box.line.fill.background()
log_text = slide.shapes.add_textbox(Inches(6.95), Inches(2.05), Inches(5.2), Inches(2.8))
tf = log_text.text_frame
log_lines = [
    "Gate controller connected and waiting for incoming messages...",
    "Summary synchronized: current free spaces = 1",
    "Summary synchronized: current free spaces = 4",
    "Gate opened: CAR-101 is allowed to enter",
    "Vehicle request sent: {'vehicle_id': 'CAR-101', 'request_time': '2026-03-26 02:48:08'}",
]
for i, line in enumerate(log_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.name = "Courier New"
    p.space_after = Pt(5)
bottom = (
    "These results confirm that the core MQTT workflow already works as intended. "
    "The prototype is simple, but it clearly demonstrates event publishing, topic subscription, shared broker communication, and decision making based on live data."
)
add_textbox(slide, 0.9, 5.35, 11.5, 1.0, [(bottom, 14, False)], fill=LIGHT, color=NAVY)


# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "6. Limitations and Future Improvements",
    "This version is a coursework prototype, so several features can still be expanded later.",
)
add_bullet_box(
    slide,
    0.85,
    1.8,
    5.7,
    4.9,
    "Current limitations",
    [
        "The project uses a public MQTT broker instead of a private and secured deployment.",
        "Parking data is simulated randomly and is not connected to real hardware sensors.",
        "The output is shown in terminal windows rather than in a graphical web dashboard.",
        "The current workflow does not store historical data in a database for later analysis.",
    ],
)
add_bullet_box(
    slide,
    6.8,
    1.8,
    5.7,
    4.9,
    "Possible future work",
    [
        "Replace the simulated sensor input with real ultrasonic or infrared sensors.",
        "Add a database to record parking events, traffic patterns, and gate decisions.",
        "Develop a web or mobile dashboard so that users can check parking availability remotely.",
        "Extend the system with reservation, pricing, and license plate recognition features.",
    ],
)


# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(
    slide,
    "7. Conclusion",
    "The project meets the coursework goal by presenting a working MQTT-based smart parking prototype.",
)
conclusion = (
    "In conclusion, this prototype shows how MQTT can be applied to a smart parking scenario in a clear and practical way. "
    "The project includes message publishers, message subscribers, a shared broker, topic design, statistics processing, and gate decision logic. "
    "Although the system is still small, it already demonstrates the essential ideas behind IoT communication and distributed services. "
    "This makes it a suitable foundation for further development in later coursework stages."
)
add_textbox(slide, 0.95, 1.95, 11.3, 2.0, [("Project summary", 21, True), (conclusion, 16, False)], fill=WHITE)
final_box = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(4.65), Inches(9.7), Inches(1.3)
)
final_box.fill.solid()
final_box.fill.fore_color.rgb = TEAL
final_box.line.fill.background()
tf = final_box.text_frame
tf.text = "Smart Parking with MQTT: a simple prototype, a clear message flow, and a strong base for future extension."
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(21)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].font.name = "Aptos"
thanks = slide.shapes.add_textbox(Inches(0.95), Inches(6.35), Inches(11.3), Inches(0.4))
thanks.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
run = thanks.text_frame.paragraphs[0].add_run()
run.text = "Thank you."
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = NAVY
run.font.name = "Aptos"

OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))
print(OUTPUT)
