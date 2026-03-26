from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/wzm/学习/工作/喷漆/code/repo_zip_download_20260326_0153/extracted/--main")
ASSETS = ROOT / "presentation" / "assets"
OUTPUT = ROOT / "presentation" / "Smart_Parking_MQTT_Presentation_EN.pptx"

BG = RGBColor(246, 248, 252)
NAVY = RGBColor(24, 44, 76)
BLUE = RGBColor(34, 92, 165)
TEAL = RGBColor(37, 136, 127)
ORANGE = RGBColor(219, 131, 58)
LIGHT = RGBColor(228, 236, 247)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(61, 72, 88)


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def top_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.42)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()


def add_title(slide, title, subtitle=""):
    top_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(12.0), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = "Aptos"
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT
        r2.font.name = "Aptos"


def add_bullet_card(slide, left, top, width, height, title, bullets, fill=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LIGHT

    box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.36), Inches(height - 0.24))
    tf = box.text_frame
    head = tf.paragraphs[0]
    head.text = title
    head.font.size = Pt(18)
    head.font.bold = True
    head.font.color.rgb = NAVY
    head.font.name = "Aptos"
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.font.name = "Aptos"
        p.space_after = Pt(4)


def add_paragraph_box(slide, left, top, width, height, title, body, fill=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LIGHT
    box = slide.shapes.add_textbox(Inches(left + 0.18), Inches(top + 0.12), Inches(width - 0.36), Inches(height - 0.24))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Aptos"
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT
    p2.font.name = "Aptos"


def add_image(slide, path, left, top, width=None, height=None):
    if width is not None and height is not None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    elif width is not None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
top_bar(slide)
hero = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(1.05), Inches(11.7), Inches(5.45)
)
hero.fill.solid()
hero.fill.fore_color.rgb = WHITE
hero.line.color.rgb = LIGHT
tb = slide.shapes.add_textbox(Inches(1.18), Inches(1.55), Inches(10.7), Inches(3.0))
tf = tb.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Smart Parking MQTT Prototype"
r.font.size = Pt(30)
r.font.bold = True
r.font.color.rgb = NAVY
r.font.name = "Aptos"
p2 = tf.add_paragraph()
r2 = p2.add_run()
r2.text = "Current Progress Presentation"
r2.font.size = Pt(18)
r2.font.color.rgb = BLUE
r2.font.name = "Aptos"
p3 = tf.add_paragraph()
r3 = p3.add_run()
r3.text = (
    "This version of the presentation focuses on the current implementation progress. "
    "Instead of repeating the project background, it shows the modules that are already working, "
    "the live demo output captured from the current codebase, and the recent GitHub development record."
)
r3.font.size = Pt(18)
r3.font.color.rgb = TEXT
r3.font.name = "Aptos"
tag = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.18), Inches(5.0), Inches(3.05), Inches(0.55)
)
tag.fill.solid()
tag.fill.fore_color.rgb = TEAL
tag.line.fill.background()
tag.text_frame.text = "Live demo + GitHub progress"
tag.text_frame.paragraphs[0].font.size = Pt(16)
tag.text_frame.paragraphs[0].font.bold = True
tag.text_frame.paragraphs[0].font.color.rgb = WHITE
tag.text_frame.paragraphs[0].font.name = "Aptos"


# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "1. Current Implementation Progress", "The core MQTT workflow is already running end to end.")
add_bullet_card(
    slide,
    0.85,
    1.8,
    4.0,
    4.9,
    "Completed modules",
    [
        "Parking sensor simulation is publishing slot updates to MQTT topics.",
        "The display module is subscribing to all parking status topics and showing a live dashboard.",
        "The statistics service is counting free and occupied spaces and publishing summary messages.",
        "The gate controller is making a real-time entry decision based on current availability.",
    ],
)
add_bullet_card(
    slide,
    4.95,
    1.8,
    4.0,
    4.9,
    "New work added in this round",
    [
        "A vehicle request publisher was added so the gate logic depends on live MQTT events.",
        "The gate controller was changed from a fixed-timer demo to a message-driven controller.",
        "The README and the project structure were cleaned up so the prototype is easier to run.",
        "An English presentation deck and progress screenshots were prepared for the demo.",
    ],
    fill=WHITE,
)
add_paragraph_box(
    slide,
    9.05,
    1.8,
    3.45,
    4.9,
    "What this means",
    "At this stage, the project is no longer only a code sketch. It already behaves like a small distributed system. "
    "Different modules publish and subscribe to the same broker, and the output now demonstrates state updates, summary aggregation, and gate decisions based on current data.",
    fill=WHITE,
)


# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "2. Live Demo: Dashboard Output", "The screenshot below comes from the current run of subscriber_display.py.")
add_image(slide, ASSETS / "demo_dashboard.png", 0.78, 1.65, width=7.2, height=4.5)
add_paragraph_box(
    slide,
    8.2,
    1.75,
    4.45,
    4.3,
    "What the screenshot shows",
    "The dashboard is already receiving live parking slot messages and updating the visible status table. "
    "In this captured state, every slot has been reported, and the system is showing both occupied and free spaces in real time. "
    "This demonstrates that the subscription side of the prototype is functioning correctly and that the parking status topic design is usable in practice.",
    fill=WHITE,
)
add_paragraph_box(
    slide,
    0.9,
    6.25,
    11.75,
    0.72,
    "Progress value",
    "This output proves that the project already has a working real-time monitoring view, even though it is currently implemented in the terminal instead of a graphical web page.",
    fill=LIGHT,
)


# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "3. Live Demo: Message Publishing and Service Interaction", "The current prototype already has active publishers and reactive service logic.")
add_image(slide, ASSETS / "demo_publishers.png", 0.82, 1.72, width=6.05, height=3.45)
add_image(slide, ASSETS / "demo_services.png", 6.95, 1.72, width=5.56, height=3.45)
add_paragraph_box(
    slide,
    0.9,
    5.35,
    12.0,
    1.15,
    "Observed result",
    "On the left, the publishers are actively sending sensor updates and vehicle entry requests. On the right, the statistics service is continuously recalculating free spaces, and the gate controller is reacting to incoming requests by opening the gate when spaces are available. "
    "This means the project already supports a full message loop rather than isolated scripts.",
    fill=WHITE,
)


# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "4. GitHub Progress and Commit Record", "The repository history now reflects the recent implementation and presentation work.")
add_image(slide, ASSETS / "github_commits_framed.png", 0.72, 1.55, width=8.3, height=5.25)
add_bullet_card(
    slide,
    9.25,
    1.75,
    3.35,
    4.6,
    "Repository progress",
    [
        "The latest branch already contains the prototype improvements and the presentation files.",
        "The commit history shows that the work was recorded step by step instead of only being kept locally.",
        "This is useful for collaboration because team members can review the exact changes made to the project.",
        "The branch can now be used as a clean checkpoint for the next round of development.",
    ],
    fill=WHITE,
)


# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide)
add_title(slide, "5. Current Status and Next Step", "The prototype is already ready for a short classroom demonstration.")
add_bullet_card(
    slide,
    0.85,
    1.85,
    5.7,
    4.8,
    "What is ready now",
    [
        "The project can run as a multi-process MQTT demo.",
        "The dashboard, statistics service, and gate logic are already producing visible output.",
        "The repository contains code updates, documentation, and presentation material.",
        "The current progress is strong enough to demonstrate a working prototype in class.",
    ],
)
add_bullet_card(
    slide,
    6.8,
    1.85,
    5.7,
    4.8,
    "Reasonable next step",
    [
        "Improve the visual interface by replacing terminal output with a simple web dashboard.",
        "Store parking events in a database so the team can analyze historical behavior later.",
        "Connect the logic to real sensors or a more realistic hardware simulation.",
        "Keep extending the repository with cleaner version control and team collaboration records.",
    ],
)


OUTPUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUTPUT))
print(OUTPUT)
